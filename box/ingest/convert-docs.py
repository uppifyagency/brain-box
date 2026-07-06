#!/usr/bin/env python3
"""Converte un documento in markdown (testo). uso: convert-docs.py <input> <output.md>
Formati: pdf docx xlsx pptx txt html htm rtf csv json epub png jpg jpeg webp.
Exit != 0 = conversione fallita (→ quarantena). Exit 3 = niente testo estraibile.
OCR (glm-ocr via Ollama) per: immagini, e pagine PDF senza text-layer (Fase 2 audit)."""
import json
import os
import sys
from pathlib import Path

MAX_XLSX_ROWS = 2000  # ponytail: un foglio enorme diventa rumore, si tronca dichiarando
MAX_JSON_CHARS = 100_000
MIN_PAGE_CHARS = 50      # sotto questa soglia una pagina PDF è "senza testo" → candidata OCR
MIN_DOC_CHARS = 200      # sotto questa soglia il DOCUMENTO è un guscio vuoto → quarantena
OCR_MAX_PAGES = int(os.environ.get("DELERA_OCR_MAX_PAGES", "10"))  # ponytail: ~9min/pagina su CPU; oltre si dichiara il troncamento
OCR_TIMEOUT_S = int(os.environ.get("DELERA_OCR_TIMEOUT_S", "900"))
OCR_MODEL = os.environ.get("DELERA_OCR_MODEL", "glm-ocr")


def ollama_url():
    # LLM_BASE_URL è .../v1 (stile openai); l'API nativa ollama sta alla radice
    base = os.environ.get("LLM_BASE_URL", "http://ollama:11434/v1")
    return base.rsplit("/v1", 1)[0]


def ocr_image_bytes(png_bytes):
    """Testo da un'immagine via glm-ocr (Ollama nativo). Errore → eccezione (→ quarantena)."""
    import base64
    import urllib.request
    req = urllib.request.Request(
        ollama_url() + "/api/generate",
        data=json.dumps({
            "model": OCR_MODEL,
            "prompt": "Extract ALL text from this image, top to bottom, without stopping early. Output the text as clean markdown, preserving tables.",
            "images": [base64.b64encode(png_bytes).decode()],
            "stream": False,
            # temperatura 0 + tetto esplicito: senza, glm-ocr a volte si ferma dopo poche righe
            "options": {"temperature": 0, "num_predict": 4096},
        }).encode(),
        headers={"Content-Type": "application/json"},
    )
    with urllib.request.urlopen(req, timeout=OCR_TIMEOUT_S) as r:
        text = json.load(r).get("response", "")
    # ponytail: glm-ocr sulle pagine sparse va in LOOP (misurato: 30 echi sul golden
    # set). Cura "taglia al primo eco": una riga lunga già vista che ricompare = da lì
    # in poi è ripetizione del modello, si scarta il resto. Ceiling: un documento con
    # la stessa riga ≥20 char legittimamente ripetuta perde la seconda parte.
    seen, out = set(), []
    for line in text.splitlines():
        key = line.strip()
        if len(key) >= 20 and key in seen:
            break
        if key:
            seen.add(key)
        out.append(line)
    return "\n".join(out)


def image_text(p):
    data = Path(p).read_bytes()
    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(data))
        # ponytail: misurato sul golden set — glm-ocr TRONCA sopra ~200dpi (A4 intera a
        # 150dpi = affidabile). Sopra 1800px di lato si ridimensiona prima dell'OCR.
        MAX_SIDE = 1800
        if max(img.size) > MAX_SIDE:
            img.thumbnail((MAX_SIDE, MAX_SIDE))
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG")
            data = buf.getvalue()
    except Exception:
        pass  # PIL assente o formato strano: si tenta l'OCR sul file com'è
    return ocr_image_bytes(data)


def pdf_text(p):
    import fitz  # PyMuPDF
    out, ocr_done, ocr_skipped = [], 0, 0
    with fitz.open(p) as doc:
        for page in doc:
            text = page.get_text("text")
            if len(text.strip()) < MIN_PAGE_CHARS:
                # pagina senza text-layer → corsia OCR (solo queste, con tetto dichiarato)
                if ocr_done < OCR_MAX_PAGES:
                    png = page.get_pixmap(dpi=150).tobytes("png")
                    text = ocr_image_bytes(png)
                    ocr_done += 1
                else:
                    ocr_skipped += 1
            out.append(text)
    if ocr_skipped:
        out.append(f"\n> (OCR troncato: {ocr_skipped} pagine senza testo oltre il tetto di {OCR_MAX_PAGES} — alza DELERA_OCR_MAX_PAGES per processarle)")
    text = "\n\n".join(out)
    # control-char grezzi dai PDF (muro storico): rompono JSON/YAML a valle
    return "".join(c if c == "\n" or c == "\t" or ord(c) >= 32 else " " for c in text)


def docx_text(p):
    from docx import Document
    doc = Document(p)
    out = [para.text for para in doc.paragraphs]
    for t in doc.tables:
        for row in t.rows:
            out.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(out)


def xlsx_text(p):
    from openpyxl import load_workbook
    wb = load_workbook(p, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"## Foglio: {ws.title}\n")
        n = 0
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append("| " + " | ".join(cells) + " |")
                n += 1
            if n >= MAX_XLSX_ROWS:
                out.append(f"\n> (troncato a {MAX_XLSX_ROWS} righe)")
                break
    return "\n".join(out)


def csv_text(p):
    import csv
    out, n = [], 0
    with open(p, newline="", errors="replace") as f:
        sniff = f.read(4096); f.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sniff, delimiters=",;\t")
        except csv.Error:
            dialect = csv.excel
        for row in csv.reader(f, dialect):
            if any(c.strip() for c in row):
                out.append("| " + " | ".join(c.strip() for c in row) + " |")
                if n == 0:
                    out.append("|" + "---|" * len(row))
                n += 1
            if n >= MAX_XLSX_ROWS:
                out.append(f"\n> (troncato a {MAX_XLSX_ROWS} righe)")
                break
    return "\n".join(out)


def json_text(p):
    data = json.loads(Path(p).read_text(errors="replace"))
    pretty = json.dumps(data, indent=2, ensure_ascii=False)
    if len(pretty) > MAX_JSON_CHARS:
        pretty = pretty[:MAX_JSON_CHARS] + f"\n… (troncato a {MAX_JSON_CHARS} caratteri)"
    return "```json\n" + pretty + "\n```"


def html_text(p):
    return _strip_html(Path(p).read_text(errors="replace"))


def _strip_html(html):
    from html.parser import HTMLParser

    class Strip(HTMLParser):
        SKIP = {"script", "style", "nav", "footer", "header"}
        def __init__(self):
            super().__init__(); self.parts = []; self.skip = 0
        def handle_starttag(self, tag, attrs):
            if tag in self.SKIP: self.skip += 1
        def handle_endtag(self, tag):
            if tag in self.SKIP and self.skip: self.skip -= 1
        def handle_data(self, d):
            if not self.skip and d.strip(): self.parts.append(d.strip())

    s = Strip()
    s.feed(html)
    return "\n".join(s.parts)


def epub_text(p):
    # ponytail: capitoli in ordine di nome dentro lo zip, non di spine OPF — per i libri
    # reali è quasi sempre equivalente; se un epub esce disordinato, si legge l'OPF.
    import zipfile
    out = []
    with zipfile.ZipFile(p) as z:
        for name in sorted(z.namelist()):
            if name.lower().endswith((".xhtml", ".html", ".htm")):
                out.append(_strip_html(z.read(name).decode(errors="replace")))
    return "\n\n".join(out)


def rtf_text(p):
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(Path(p).read_text(errors="replace"))


def pptx_text(p):
    from pptx import Presentation
    prs = Presentation(p)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.extend(par.text for par in shape.text_frame.paragraphs if par.text.strip())
    return "\n".join(out)


CONVERT = {
    ".pdf": pdf_text, ".docx": docx_text, ".xlsx": xlsx_text, ".pptx": pptx_text,
    ".html": html_text, ".htm": html_text, ".rtf": rtf_text,
    ".txt": lambda p: Path(p).read_text(errors="replace"),
    ".text": lambda p: Path(p).read_text(errors="replace"),
    ".csv": csv_text, ".json": json_text, ".epub": epub_text,
    ".png": image_text, ".jpg": image_text, ".jpeg": image_text, ".webp": image_text,
}


# Gate anti-guscio: il falso successo da beccare è "file SOSTANZIOSO da cui esce quasi
# nulla" (estrazione fallita). Lezioni del golden set: (1) sui formati verbatim
# (txt/csv/json) un file corto è contenuto legittimo — niente gate; (2) anche docx/pdf
# corti ma PICCOLI sono legittimi → il gate scatta solo se il sorgente è grande;
# (3) per le immagini l'OCR restituisce quel che c'è: basta il non-vuoto.
EXTRACTED_FORMATS = {".pdf", ".docx", ".pptx", ".xlsx", ".epub"}
MIN_SRC_BYTES_FOR_GATE = 50_000


def main():
    src, dst = Path(sys.argv[1]), Path(sys.argv[2])
    ext = src.suffix.lower()
    fn = CONVERT.get(ext)
    if fn is None:
        print(f"formato non supportato: {src.suffix}", file=sys.stderr); return 2
    text = fn(str(src)).strip()
    if not text:
        print("documento vuoto dopo l'estrazione", file=sys.stderr); return 3
    if (ext in EXTRACTED_FORMATS and len(text) < MIN_DOC_CHARS
            and src.stat().st_size > MIN_SRC_BYTES_FOR_GATE):
        print(f"niente testo estraibile ({len(text)} caratteri da {src.stat().st_size} byte) — estrazione probabilmente fallita", file=sys.stderr)
        return 3
    dst.write_text(text + "\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
